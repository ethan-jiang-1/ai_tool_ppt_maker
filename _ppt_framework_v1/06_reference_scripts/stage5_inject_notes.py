#!/usr/bin/env python3
"""
Stage 5: Inject speaker notes into the PPTX from the source markdown.

Extracts SPEAKER NOTE blocks from the same markdown files used in Stage 1,
injects each note into the corresponding slide's notes panel in the PPTX.

⚠️  This stage modifies the PPTX in-place. Back it up first.

Usage:
    python stage5_inject_notes.py \\
        --pptx 3_versions/v1/_generated/ppt/{NAME}.pptx \\
        --input 3_versions/v1/slide-specifications.md

Dependencies: python-pptx
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation


def extract_notes_from_markdown(md_paths: list[str]) -> list[str]:
    """Extract SPEAKER NOTE text from each slide block, in order."""
    all_notes: list[str] = []

    for md_path in md_paths:
        text = Path(md_path).read_text(encoding="utf-8")
        # Accept colon (:), full-width colon (：), hyphen (-), or em-dash (—) after slide number
        blocks = re.split(r"^## Slide \d+\s*[：:\-–—]", text, flags=re.M)[1:]

        for block in blocks:
            # Two formats are used across templates:
            # Format A (multi-line, pitch-deck): > **SPEAKER NOTE**\n> content...
            # Format B (one-line, others):       > **SPEAKER NOTE**: content
            note = ""

            # Try Format B first: colon + inline text on same line
            m = re.search(
                r"> \*\*SPEAKER NOTE\*\*:\s*(.+)$",
                block, re.M,
            )
            if m:
                note = m.group(1).strip()
            else:
                # Try Format A: SPEAKER NOTE on its own line, blockquote content follows
                m = re.search(
                    r"> \*\*SPEAKER NOTE\*\*\s*\n((?:> .+\n?)+)",
                    block, re.M,
                )
                if m:
                    note = re.sub(r"^> ?", "", m.group(1), flags=re.M).strip()

            all_notes.append(note)

    return all_notes


def main() -> None:
    parser = argparse.ArgumentParser(description="Stage 5: Inject speaker notes")
    parser.add_argument("--pptx", required=True, help="PPTX file to modify in-place")
    parser.add_argument("--input", nargs="+", required=True,
                        help="One or more markdown slide spec files")
    args = parser.parse_args()

    notes = extract_notes_from_markdown(args.input)
    prs = Presentation(args.pptx)

    # Notes are matched to slides positionally, so the counts MUST agree. They only
    # diverge when the deck was built partial (Stage 3/4 now fail loud to prevent
    # that) or the spec was edited after the .pptx was built. Aborting here beats
    # the old warn-and-continue, which silently shifted every note after the gap.
    if len(notes) != len(prs.slides):
        raise SystemExit(
            f"✗ Stage 5 aborted: {len(notes)} speaker notes in the spec but "
            f"{len(prs.slides)} slides in the PPTX. Positional matching would misalign "
            f"notes. Rebuild the deck from the current spec (Stage 1→4) so counts agree, "
            f"then rerun Stage 5.")

    count = 0
    for slide, note_text in zip(prs.slides, notes):
        if note_text:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            tf.text = note_text
            count += 1

    prs.save(args.pptx)

    print(f"\n--- Stage 5 complete ---")
    print(f"Notes injected: {count}/{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
