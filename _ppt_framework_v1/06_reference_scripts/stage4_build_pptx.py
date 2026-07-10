#!/usr/bin/env python3
"""
Stage 4: Build the PPTX container — wrap final slide images into a .pptx file.

Reads header-locked PNG images from Stage 3 and slide_plan.json from Stage 1.
Creates a 16:9 PPTX with one full-bleed image per slide. No editable text objects
— the PPTX is a media container; all content is in the images.

Usage:
    python stage4_build_pptx.py \\
        --images header_locked/ \\
        --slide-plan slide_plan.json \\
        --out ppt/deck.pptx \\
        --title "My Presentation"

Dependencies: python-pptx
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


# 16:9 standard
SLIDE_WIDTH_IN = Inches(13.333)
SLIDE_HEIGHT_IN = Inches(7.5)

_IMG_EXTS = {".png", ".jpg", ".jpeg"}


def match_slide_image(img_dir: Path, slide_id: str) -> list[Path]:
    """Images matching a slide id under the canonical NN_<id> (or bare <id>) naming
    — ANCHORED and sorted. An unanchored substring match cross-hits ids like 's1'
    onto '10_s10.png'; anchoring the id to the stem's end prevents wrong-image builds."""
    pat = re.compile(rf"^(\d+_)?{re.escape(slide_id)}$")
    return sorted(p for p in img_dir.iterdir()
                  if p.is_file() and p.suffix.lower() in _IMG_EXTS and pat.match(p.stem))


def main() -> None:
    parser = argparse.ArgumentParser(description="Stage 4: Build PPTX container")
    parser.add_argument("--images", required=True,
                        help="Directory containing Stage 3 header-locked images")
    parser.add_argument("--slide-plan", required=True,
                        help="slide_plan.json from Stage 1")
    parser.add_argument("--out", required=True, help="Output .pptx path")
    parser.add_argument("--title", default="Presentation", help="Deck title")
    args = parser.parse_args()

    plan_data = json.loads(Path(args.slide_plan).read_text(encoding="utf-8"))
    slides = plan_data.get("slides", [])

    img_dir = Path(args.images)

    # Fail loud on a partial/ambiguous image set BEFORE building — a skipped slide
    # would silently shrink the deck and misalign downstream speaker notes.
    resolved: list[Path] = []
    problems: list[str] = []
    for slide in slides:
        sid = slide["id"]
        hits = match_slide_image(img_dir, sid)
        if not hits:
            problems.append(f"no image for slide {sid!r} (expected NN_{sid}.png in {img_dir.name}/)")
        elif len(hits) > 1:
            problems.append(f"ambiguous images for slide {sid!r}: {[h.name for h in hits]}")
        else:
            resolved.append(hits[0])
    if problems:
        raise SystemExit(
            f"✗ Stage 4 cannot build the deck — {len(problems)} image problem(s):\n" +
            "\n".join(f"  - {p}" for p in problems) +
            f"\n  Re-run Stage 3 (and Stage 2 if images are missing) so every slide has "
            f"exactly one header-locked image, then Stage 4.")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_IN
    prs.slide_height = SLIDE_HEIGHT_IN

    blank_layout = prs.slide_layouts[6]  # blank — no placeholders

    for i, (slide, img_path) in enumerate(zip(slides, resolved)):
        slide_id = slide["id"]
        new_slide = prs.slides.add_slide(blank_layout)
        new_slide.shapes.add_picture(
            str(img_path),
            left=Inches(0),
            top=Inches(0),
            width=SLIDE_WIDTH_IN,
            height=SLIDE_HEIGHT_IN,
        )
        print(f"  Slide {i + 1}: {slide_id}  ← {img_path.name}")

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    print(f"\n--- Stage 4 complete ---")
    print(f"PPTX: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
