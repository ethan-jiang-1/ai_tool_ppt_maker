#!/usr/bin/env python3
"""Offline end-to-end smoke test for Stages 1, 3, 4, and 5."""

from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))
import bundle_layout as layout

HERE = Path(__file__).resolve().parent


def test_offline_pipeline_builds_pptx_and_notes():
    with tempfile.TemporaryDirectory() as td:
        deck = Path(td) / "deck_e2e"
        layout.init_bundle(deck, style="dark-executive")
        run_dir = deck / layout.VERSIONS_DIR / "v1"
        spec = run_dir / layout.SLIDE_SPECS_NAME
        spec.write_text(
            """## Slide 1: intro

**VISUAL TYPE**: Data / Chart
**RENDER MODE**: body+header-lock
**KICKER**: CONTEXT
**TITLE**: The pipeline is reproducible

**IMAGE PROMPT**:
```
A single clean proof card on a dark navy background.
```

> **SPEAKER NOTE**: This note proves Stage 5 stayed aligned.
""",
            encoding="utf-8")

        stage1 = subprocess.run(
            [sys.executable, str(HERE / "unified_pipeline.py"),
             "--run-dir", str(run_dir), "--stage", "1"],
            capture_output=True, text=True)
        assert stage1.returncode == 0, stage1.stdout + stage1.stderr

        images = layout.generated_dir(run_dir) / layout.GEN_IMAGES_SUBDIR
        images.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (1672, 941), (10, 22, 42)).save(images / "01_intro.png")

        finish = subprocess.run(
            [sys.executable, str(HERE / "unified_pipeline.py"),
             "--run-dir", str(run_dir), "--stage", "3,4,5"],
            capture_output=True, text=True)
        assert finish.returncode == 0, finish.stdout + finish.stderr

        pptx = layout.generated_dir(run_dir) / layout.GEN_PPT_SUBDIR / "e2e.pptx"
        backup = pptx.with_suffix(".backup.pptx")
        assert pptx.is_file() and backup.is_file()
        prs = Presentation(pptx)
        assert len(prs.slides) == 1
        assert "Stage 5 stayed aligned" in prs.slides[0].notes_slide.notes_text_frame.text


def _main() -> int:
    try:
        test_offline_pipeline_builds_pptx_and_notes()
        print("  ✓ test_offline_pipeline_builds_pptx_and_notes")
        print("\n✓ all passed (1 test)")
        return 0
    except Exception as exc:
        print(f"  ✗ test_offline_pipeline_builds_pptx_and_notes ({type(exc).__name__}: {exc})")
        return 1


if __name__ == "__main__":
    raise SystemExit(_main())
